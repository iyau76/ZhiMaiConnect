"""Generate the editable five-minute Zhimai Connect competition deck."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "doc" / "知脉Connect-5分钟演示稿-待填队长院系.pptx"
LOGO = ROOT / "src" / "assets" / "logo-mark-384.png"

FONT_CN = "Microsoft YaHei"
FONT_LATIN = "Aptos"

INK = RGBColor(31, 31, 45)
MUTED = RGBColor(96, 92, 111)
PAPER = RGBColor(249, 248, 252)
WHITE = RGBColor(255, 255, 255)
VIOLET = RGBColor(112, 31, 235)
PINK = RGBColor(246, 79, 139)
CORAL = RGBColor(255, 91, 92)
CYAN = RGBColor(39, 181, 176)
LILAC = RGBColor(235, 224, 255)
ROSE = RGBColor(255, 229, 239)
MINT = RGBColor(221, 246, 241)
LINE = RGBColor(219, 215, 228)
SOFT = RGBColor(242, 239, 247)
DARK = RGBColor(22, 20, 35)

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)


def set_font(run, *, size=18, bold=False, color=INK, latin=False):
    run.font.name = FONT_LATIN if latin else FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), FONT_CN)
    rpr.set(qn("a:latin"), FONT_LATIN)


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    line_spacing=1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return box


def add_rich_lines(slide, x, y, w, h, lines, *, fill=WHITE, border=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.22)
    frame.margin_top = frame.margin_bottom = Inches(0.15)
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(line.get("after", 5))
        paragraph.level = line.get("level", 0)
        paragraph.alignment = line.get("align", PP_ALIGN.LEFT)
        run = paragraph.add_run()
        run.text = line["text"]
        set_font(
            run,
            size=line.get("size", 15),
            bold=line.get("bold", False),
            color=line.get("color", INK),
        )
    return shape


def add_pill(slide, x, y, w, text, *, fill=SOFT, color=INK, border=None, size=11):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border or fill
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.08)
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    set_font(run, size=size, bold=True, color=color)
    return shape


def add_circle_label(slide, x, y, diameter, text, *, fill=VIOLET, color=WHITE, size=16):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    set_font(run, size=size, bold=True, color=color)
    return shape


def add_background(slide, color=PAPER):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = color


def add_header(slide, number, section, title, subtitle, seconds):
    add_text(slide, 0.62, 0.34, 3.2, 0.25, section.upper(), size=9, color=VIOLET, bold=True)
    add_text(slide, 0.62, 0.69, 11.7, 0.62, title, size=27, color=INK, bold=True)
    if subtitle:
        add_text(slide, 0.64, 1.28, 11.6, 0.38, subtitle, size=12, color=MUTED)
    add_pill(slide, 11.91, 0.32, 0.78, f"{seconds} 秒", fill=LILAC, color=VIOLET, size=10)
    add_text(slide, 12.12, 7.10, 0.55, 0.2, f"{number:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, text="仅使用合成演示数据 · 不抓取社交平台 · 不自动发送消息"):
    add_text(slide, 0.64, 7.08, 10.8, 0.18, text, size=8, color=MUTED)


def add_arrow(slide, x1, y1, x2, y2, color=VIOLET, width=2.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK)
    # Accent rail and a generous white logo tile keep the raster mark crisp.
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), SLIDE_H)
    rail.fill.solid()
    rail.fill.fore_color.rgb = PINK
    rail.line.fill.background()
    tile = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.72), Inches(0.72), Inches(2.68), Inches(2.68)
    )
    tile.fill.solid()
    tile.fill.fore_color.rgb = WHITE
    tile.line.color.rgb = WHITE
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(9.91), Inches(0.91), Inches(2.30), Inches(2.30))
    add_pill(slide, 0.72, 0.72, 3.18, "LOCAL-FIRST · EVIDENCE-BOUND", fill=VIOLET, color=WHITE, size=10)
    add_text(slide, 0.72, 1.52, 8.6, 0.95, "知脉 Connect", size=42, color=WHITE, bold=True)
    add_text(
        slide,
        0.75,
        2.55,
        8.2,
        1.10,
        "把零散关系记忆，变成\n可复核、可追溯、可行动的个人知识",
        size=24,
        color=RGBColor(235, 231, 244),
        bold=True,
        line_spacing=1.0,
    )
    add_pill(slide, 0.75, 4.12, 1.55, "用户主动提供", fill=RGBColor(50, 43, 72), color=WHITE)
    add_pill(slide, 2.48, 4.12, 1.55, "确认后入库", fill=RGBColor(50, 43, 72), color=WHITE)
    add_pill(slide, 4.21, 4.12, 1.55, "本地确定排序", fill=RGBColor(50, 43, 72), color=WHITE)
    add_pill(slide, 5.94, 4.12, 1.75, "AI 只做解释润色", fill=RGBColor(50, 43, 72), color=WHITE)
    add_text(
        slide,
        0.75,
        5.33,
        8.9,
        0.42,
        "Vibe Coding 赛道 · 五分钟产品演示",
        size=14,
        color=RGBColor(190, 183, 205),
    )
    add_text(
        slide,
        0.75,
        6.34,
        9.5,
        0.35,
        "队长：[待填写]   ·   院系：[待填写]   ·   部署二维码：[待填写]",
        size=11,
        color=RGBColor(255, 177, 210),
        bold=True,
    )
    add_pill(slide, 11.64, 6.70, 0.92, "15 秒", fill=PINK, color=WHITE, size=10)
    return slide


def pain_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, 2, "01 · PAIN", "真正的断点发生在“要行动”的那一刻", "关系信息并不稀缺，但缺少可信的调用方式。", 25)
    add_text(slide, 0.68, 1.88, 3.1, 0.35, "碎片化输入", size=14, color=MUTED, bold=True)
    chips = [("聊天截图", ROSE), ("活动名单", LILAC), ("语音记忆", MINT), ("脑海印象", SOFT)]
    for index, (label, color) in enumerate(chips):
        add_pill(slide, 0.70 + (index % 2) * 1.48, 2.35 + (index // 2) * 0.62, 1.30, label, fill=color)
    add_arrow(slide, 3.80, 2.85, 4.82, 2.85, PINK, 3)
    questions = [
        ("1", "该找谁？", "不是永久排名，而是当前任务最合适的人"),
        ("2", "依据是什么？", "技能、共同事件、来源与风险必须可查"),
        ("3", "信息还有效吗？", "时间、状态与待确认不能被模型补写"),
    ]
    for index, (num, title, body) in enumerate(questions):
        y = 1.92 + index * 1.22
        add_circle_label(slide, 5.02, y, 0.48, num, fill=[VIOLET, PINK, CYAN][index], size=10)
        add_rich_lines(
            slide,
            5.66,
            y - 0.03,
            6.72,
            0.82,
            [
                {"text": title, "size": 17, "bold": True},
                {"text": body, "size": 11, "color": MUTED, "after": 0},
            ],
        )
    add_rich_lines(
        slide,
        0.70,
        5.84,
        11.70,
        0.72,
        [
            {"text": "产品判断", "size": 11, "bold": True, "color": VIOLET},
            {"text": "关系图只是界面；真正的价值是把记忆变成有证据、可确认的下一步行动。", "size": 17, "bold": True},
        ],
        fill=WHITE,
        border=LILAC,
    )
    add_footer(slide)


def users_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, 3, "02 · USERS & BOUNDARIES", "服务关系密集型人群，也主动限制产品权力", "先说清楚为谁做、什么绝对不做。", 25)
    add_rich_lines(
        slide,
        0.70,
        1.90,
        5.72,
        3.95,
        [
            {"text": "目标用户", "size": 15, "bold": True, "color": VIOLET},
            {"text": "学生组织者 / 社群负责人", "size": 19, "bold": True},
            {"text": "组织活动时，需要快速调用合作与技能记忆", "size": 11, "color": MUTED},
            {"text": "研究人员 / 自由职业者", "size": 19, "bold": True},
            {"text": "联系人多、合作周期长，信息容易过期", "size": 11, "color": MUTED},
            {"text": "招聘 / 销售等关系密集角色", "size": 19, "bold": True},
            {"text": "需要提醒与线索，不需要自动替人沟通", "size": 11, "color": MUTED},
        ],
        fill=WHITE,
        border=LILAC,
    )
    add_rich_lines(
        slide,
        6.72,
        1.90,
        5.68,
        3.95,
        [
            {"text": "明确边界", "size": 15, "bold": True, "color": PINK},
            {"text": "× 不登录、抓取或逆向读取个人社交平台", "size": 16, "bold": True},
            {"text": "只处理用户主动提供的文字、文件、图片和语音", "size": 11, "color": MUTED},
            {"text": "× 不自动发送消息或修改好友状态", "size": 16, "bold": True},
            {"text": "话术只能编辑、复制，由用户决定是否发送", "size": 11, "color": MUTED},
            {"text": "× 不把模型推断当作事实", "size": 16, "bold": True},
            {"text": "AI 输出先进入草稿，仍有待确认就不能入库", "size": 11, "color": MUTED},
        ],
        fill=WHITE,
        border=ROSE,
    )
    add_pill(slide, 3.38, 6.17, 6.58, "隐私不是页脚文案，而是产品流程中的权限边界", fill=DARK, color=WHITE, size=13)
    add_footer(slide)


def loop_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, 4, "03 · PRODUCT LOOP", "四步闭环：从主动材料到可解释行动", "每一步都保留用户控制权。", 35)
    cards = [
        ("1", "低摩擦录入", "文字 / 文件 / 截图 / 语音", LILAC, VIOLET),
        ("2", "证据门禁 + 复核", "清空无证据值；逐项接受或拒绝", ROSE, PINK),
        ("3", "本地关系记忆", "人物 / 身份 / 关系 / 事件 / 提醒", MINT, CYAN),
        ("4", "行动建议", "确定性 Top 3 + 可编辑 AI 解释", SOFT, INK),
    ]
    for index, (num, title, body, fill, accent) in enumerate(cards):
        x = 0.70 + index * 3.05
        add_rich_lines(
            slide,
            x,
            2.02,
            2.64,
            2.65,
            [
                {"text": num, "size": 28, "bold": True, "color": accent},
                {"text": title, "size": 18, "bold": True},
                {"text": body, "size": 12, "color": MUTED},
            ],
            fill=fill,
            border=fill,
        )
        if index < 3:
            add_arrow(slide, x + 2.67, 3.34, x + 3.00, 3.34, accent, 2)
    add_rich_lines(
        slide,
        0.70,
        5.18,
        11.80,
        1.02,
        [
            {"text": "外部模型不可用时仍可演示", "size": 13, "bold": True, "color": VIOLET},
            {"text": "IndexedDB 本地资料、关系图、日期提醒与候选排序继续工作；AI 只负责可选的比较说明和话术润色。", "size": 15, "bold": True},
        ],
        fill=WHITE,
        border=LINE,
    )
    add_footer(slide)


def demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, 5, "04 · LIVE DEMO", "一条合成故事，连续验证三个场景", "固定故事：筹备“校园记忆展”；所有人物与数据均为虚构。", 95)
    demos = [
        (
            "A",
            "录入与确认",
            "粘贴材料 → AI 草稿\n编辑一项、拒绝一项\n身份历史与关系入库",
            VIOLET,
            LILAC,
            "45 秒",
        ),
        (
            "B",
            "这事找谁",
            "先本地稳定召回 Top 3\n展示证据、时间与风险\n再生成比较与求助话术",
            PINK,
            ROSE,
            "30 秒",
        ),
        (
            "C",
            "提醒与维护",
            "生日 / 节日 / 长期未联系\n资料不足时明确提示缺口\n日历回看共同事件",
            CYAN,
            MINT,
            "20 秒",
        ),
    ]
    for index, (letter, title, body, accent, fill, timing) in enumerate(demos):
        x = 0.70 + index * 4.02
        add_rich_lines(
            slide,
            x,
            1.95,
            3.60,
            3.85,
            [
                {"text": f"场景 {letter}", "size": 12, "bold": True, "color": accent},
                {"text": title, "size": 22, "bold": True},
                {"text": body, "size": 13, "color": INK, "after": 8},
                {"text": timing, "size": 11, "bold": True, "color": accent},
            ],
            fill=fill,
            border=fill,
        )
    add_pill(slide, 3.18, 6.16, 7.00, "演示原则：先展示本地确定结果，再触发可选 AI；全程不外发消息", fill=DARK, color=WHITE, size=12)
    add_footer(slide)


def architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, 6, "05 · ARCHITECTURE", "难点不在调用模型，而在可验证的数据契约", "React/TanStack + IndexedDB；AI 路由是受限的可选能力。", 35)
    layers = [
        ("输入层", "文字 · 文件 · 图片 · 语音", LILAC, VIOLET),
        ("复核层", "Schema · 主体分段 · 正向绑定 · 否定拦截", ROSE, PINK),
        ("本地数据层", "Person · Identity · Relation · Event · Fact", MINT, CYAN),
        ("行动层", "局部图谱 · 提醒 · Top 3 · 风险解释", SOFT, INK),
    ]
    for index, (title, body, fill, accent) in enumerate(layers):
        y = 1.85 + index * 1.04
        add_rich_lines(
            slide,
            0.72,
            y,
            7.25,
            0.78,
            [
                {"text": title, "size": 14, "bold": True, "color": accent},
                {"text": body, "size": 12, "color": INK},
            ],
            fill=fill,
            border=fill,
        )
        if index < 3:
            add_arrow(slide, 4.33, y + 0.80, 4.33, y + 1.01, accent, 1.5)
    add_rich_lines(
        slide,
        8.35,
        1.86,
        4.02,
        2.05,
        [
            {"text": "可选 AI 通道", "size": 16, "bold": True, "color": VIOLET},
            {"text": "首次按服务商与数据类型确认", "size": 12},
            {"text": "只发送当前任务所需内容", "size": 12},
            {"text": "AI 不改变本地候选排序", "size": 12},
        ],
        fill=WHITE,
        border=LILAC,
    )
    add_rich_lines(
        slide,
        8.35,
        4.18,
        4.02,
        1.84,
        [
            {"text": "公网路由最小防护", "size": 16, "bold": True, "color": PINK},
            {"text": "会话握手 · 限速 · 大小限制", "size": 12},
            {"text": "超时 · SSRF · no-store · 安全错误", "size": 12},
        ],
        fill=WHITE,
        border=ROSE,
    )
    add_footer(slide)


def innovation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, 7, "06 · DIFFERENTIATION", "四项差异化，共同服务“可信行动”", "不是更会猜，而是更清楚何时不能猜。", 30)
    items = [
        ("01", "身份是历史", "平台账号、历史昵称与有效期，降低改名失联和同名误合并", VIOLET, LILAC),
        ("02", "事实有证据", "来源、时间、确认状态与风险可见；无证据的 AI 值先清空", PINK, ROSE),
        ("03", "推荐可复现", "本地规则召回 Top 3，AI 解释“为什么是他 / 不是另一个人”", CYAN, MINT),
        ("04", "采集最小暴露", "主动分享、确认入库、按次云同意；不抓取、不自动发送", INK, SOFT),
    ]
    for index, (num, title, body, accent, fill) in enumerate(items):
        x = 0.70 + (index % 2) * 6.02
        y = 1.88 + (index // 2) * 2.20
        add_rich_lines(
            slide,
            x,
            y,
            5.68,
            1.83,
            [
                {"text": num, "size": 12, "bold": True, "color": accent},
                {"text": title, "size": 20, "bold": True},
                {"text": body, "size": 12, "color": MUTED},
            ],
            fill=fill,
            border=fill,
        )
    add_pill(slide, 3.30, 6.30, 6.72, "可信度来自数据契约与用户控制，不来自模型口头保证", fill=DARK, color=WHITE, size=12)
    add_footer(slide)


def validation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, 8, "07 · VALIDATION", "演示规模与自动化门禁已经固定", "这些数字来自当前工作树的本地实测，不是预测值。", 25)
    metrics = [
        ("172", "单元测试", "18 个测试文件", VIOLET),
        ("17 × 2", "真实浏览器 E2E", "Chrome + Edge", PINK),
        ("50 / 80", "合成人物 / 关系", "一键加载、可反复重置", CYAN),
        ("3", "响应式宽度", "390 / 768 / 1440 px", INK),
    ]
    for index, (value, label, note, accent) in enumerate(metrics):
        x = 0.70 + index * 3.05
        add_rich_lines(
            slide,
            x,
            1.96,
            2.64,
            2.08,
            [
                {"text": value, "size": 29, "bold": True, "color": accent},
                {"text": label, "size": 14, "bold": True},
                {"text": note, "size": 10, "color": MUTED},
            ],
            fill=WHITE,
            border=LINE,
        )
    add_rich_lines(
        slide,
        0.70,
        4.35,
        7.55,
        1.55,
        [
            {"text": "全量门禁", "size": 13, "bold": True, "color": VIOLET},
            {"text": "typecheck · lint 0 warning · format · build · audit 0 vulnerabilities", "size": 16, "bold": True},
            {"text": "E2E 阻断外部网络，并对未知 /api/* 请求 fail-closed。", "size": 11, "color": MUTED},
        ],
        fill=LILAC,
        border=LILAC,
    )
    add_rich_lines(
        slide,
        8.55,
        4.35,
        3.84,
        1.55,
        [
            {"text": "已知边界", "size": 13, "bold": True, "color": PINK},
            {"text": "当前证据绑定仍是保守启发式；跨段代词、语义蕴含与完整证据区间列入 Track B。", "size": 12, "color": INK},
        ],
        fill=ROSE,
        border=ROSE,
    )
    add_footer(slide)


def roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK)
    add_text(slide, 0.72, 0.48, 2.8, 0.24, "08 · ROADMAP", size=9, color=RGBColor(190, 183, 205), bold=True)
    add_text(slide, 0.72, 0.95, 10.9, 0.72, "下一步：从可演示原型走向可托付的个人工具", size=28, color=WHITE, bold=True)
    add_pill(slide, 11.70, 0.48, 0.86, "15 秒", fill=PINK, color=WHITE, size=10)
    add_rich_lines(
        slide,
        0.72,
        2.05,
        3.65,
        3.38,
        [
            {"text": "初评前", "size": 14, "bold": True, "color": PINK},
            {"text": "部署链接与无痕复测", "size": 17, "bold": True, "color": WHITE},
            {"text": "录屏、字幕、信息表与合规自查", "size": 12, "color": RGBColor(205, 199, 218)},
            {"text": "队长完成文件命名与邮件回执", "size": 12, "color": RGBColor(205, 199, 218)},
        ],
        fill=RGBColor(41, 36, 57),
        border=RGBColor(70, 61, 91),
    )
    add_rich_lines(
        slide,
        4.62,
        2.05,
        4.34,
        3.38,
        [
            {"text": "Track B", "size": 14, "bold": True, "color": CYAN},
            {"text": "证据 ID · 原文区间 · 实体归属", "size": 15, "bold": True, "color": WHITE},
            {"text": "加密备份、恢复、审计与数据迁移", "size": 12, "color": RGBColor(205, 199, 218)},
            {"text": "15 分钟导入 30 人的真实用户试验", "size": 12, "color": RGBColor(205, 199, 218)},
        ],
        fill=RGBColor(41, 36, 57),
        border=RGBColor(70, 61, 91),
    )
    qr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.30), Inches(2.05), Inches(3.05), Inches(3.38)
    )
    qr.fill.solid()
    qr.fill.fore_color.rgb = WHITE
    qr.line.color.rgb = PINK
    qr.line.width = Pt(2)
    add_text(slide, 9.58, 2.78, 2.50, 0.55, "部署二维码", size=19, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 9.58, 3.45, 2.50, 0.58, "[待部署后替换]", size=12, color=PINK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.74, 6.12, 11.55, 0.55, "知脉 Connect：把“记得这个人”，变成“知道下一步为什么这样做”。", size=20, color=WHITE, bold=True)
    add_text(slide, 0.74, 6.85, 8.6, 0.22, "队长：[待填写] · 院系：[待填写] · 链接：[待填写]", size=10, color=RGBColor(255, 177, 210), bold=True)


def build_deck(output: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    title_slide(prs)
    pain_slide(prs)
    users_slide(prs)
    loop_slide(prs)
    demo_slide(prs)
    architecture_slide(prs)
    innovation_slide(prs)
    validation_slide(prs)
    roadmap_slide(prs)
    prs.core_properties.title = "知脉 Connect · 五分钟竞赛演示"
    prs.core_properties.subject = "本地优先、证据可追溯的人际关系记忆与行动助手"
    prs.core_properties.author = "知脉 Connect 团队（待补充队长与院系）"
    prs.core_properties.keywords = "知脉 Connect, Vibe Coding, 人际关系, 本地优先, 证据可追溯"
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    output = args.output.resolve()
    build_deck(output)
    print(output)


if __name__ == "__main__":
    main()
